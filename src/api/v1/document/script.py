from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import pytesseract
from PIL import Image
from pdf2image import convert_from_path


class DocumentParser:

    def extract_to_markdown(self, file_path: str) -> str:
        """
        Extracts text from various document types and returns it in markdown format.
        Supported formats: PDF, DOCX, PPTX, TXT, and images (PNG, JPEG).
        """

        file_path = file_path.lower()

        if file_path.endswith(".pdf"):
            return self.parse_pdf(file_path)

        elif file_path.endswith(".docx"):
            return self.parse_docx(file_path)

        elif file_path.endswith(".pptx"):
            return self.parse_pptx(file_path)

        elif file_path.endswith(".txt"):
            return self.parse_txt(file_path)

        elif file_path.endswith((".png", ".jpg", ".jpeg")):
            return self.parse_image(file_path)

        else:
            raise ValueError("Unsupported file type")

    def parse_pdf(self, file_path):

        reader = PdfReader(file_path)
        text = []

        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()

            if page_text:
                text.append(f"## Page {i+1}\n\n{page_text}")

        full_text = "\n\n".join(text)

        if not full_text.strip():
            return self.ocr_pdf(file_path)

        return full_text

    def ocr_pdf(self, file_path):

        images = convert_from_path(file_path)
        text = []

        for i, img in enumerate(images):
            ocr_text = pytesseract.image_to_string(img)
            text.append(f"## Page {i+1}\n\n{ocr_text}")

        return "\n\n".join(text)

    def parse_docx(self, file_path):

        doc = Document(file_path)
        md = []

        for para in doc.paragraphs:

            style = para.style.name.lower()

            if "heading" in style:
                level = style.replace("heading", "").strip()
                level = int(level) if level.isdigit() else 1
                md.append("#" * level + " " + para.text)

            else:
                md.append(para.text)

        return "\n\n".join(md)

    def parse_pptx(self, file_path):

        prs = Presentation(file_path)
        md = []

        for i, slide in enumerate(prs.slides):

            md.append(f"# Slide {i+1}")

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    md.append(shape.text)

        return "\n\n".join(md)

    def parse_txt(self, file_path):

        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    def parse_image(self, file_path):

        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)

        return f"# OCR Result\n\n{text}"
    

parser = DocumentParser()